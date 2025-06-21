import fitz
import docx
import os
import time
import requests
import json
import re
import logging
import datetime
from pptx import Presentation  
from dotenv import load_dotenv
from flask import Flask, request, jsonify
from flask_cors import CORS
from werkzeug.utils import secure_filename
from pymongo import MongoClient
from auth import auth

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API_KEY")

app = Flask(__name__)
CORS(app) 

app.register_blueprint(auth, url_prefix='/api/auth')

UPLOAD_FOLDER = 'uploads'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

MONGO_URI = os.getenv("MONGO_URI", "mongodb://localhost:27017")
client = MongoClient(MONGO_URI)
mongo_db = client.PassionInfotech  
history_collection = mongo_db.AI_RISK  # This collection is only for risk history

def split_text(text, chunk_size=4000):
    """Split text into smaller chunks for API processing."""
    return [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]

def extract_text_from_pdf(pdf_path):
    """Extract text content from PDF files."""
    try:
        doc = fitz.open(pdf_path)
        text = "\n".join([page.get_text("text") for page in doc])
        doc.close()
        return text
    except Exception as e:
        logger.error(f"PDF extraction error: {e}")
        return None

def extract_text_from_docx(docx_path):
    """Extract text content from DOCX files."""
    try:
        doc = docx.Document(docx_path)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        logger.error(f"DOCX extraction error: {e}")
        return None

def extract_text_from_txt(txt_path):
    """Extract text content from TXT files."""
    try:
        with open(txt_path, "r", encoding="utf-8") as file:
            return file.read()
    except Exception as e:
        logger.error(f"TXT extraction error: {e}")
        return None

def extract_text_from_pptx(pptx_path):
    """Extract text content from PPTX files."""
    try:
        presentation = Presentation(pptx_path)
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        logger.error(f"PPTX extraction error: {e}")
        return None

# Risk analysis functions
def analyze_risks_with_groq(text):
    """Process document text and generate risk reports using Groq API."""
    headers = {
        "Authorization": f"Bearer {GROQ_API_KEY}",
        "Content-Type": "application/json"
    }

    chunks = split_text(text)
    risk_reports = []
    
    for idx, chunk in enumerate(chunks):
        logger.info(f"Processing chunk {idx + 1}/{len(chunks)}...")
        
        retry_attempts = 5
        delay = 2  # Start with 2 seconds delay

        while retry_attempts > 0:
            prompt = f"""
            You are an AI specializing in risk assessment.
            Given the following document section, analyze potential risks and return ONLY properly formatted JSON.
            
            IMPORTANT FORMATTING INSTRUCTIONS:
            1. Your response must contain ONLY a single valid JSON object
            2. Do not include any explanatory text before or after the JSON
            3. Do not use markdown code blocks or triple backticks (```)
            4. Make sure all keys and string values use double quotes, not single quotes
            5. Make sure the JSON syntax is valid - test it carefully
            
            Use exactly this JSON structure:
            {{
                "RiskID": "RISK-{idx+1:03d}",
                "RiskName": "Brief name of the risk",
                "RiskCategory": "Category such as security, compliance, feasibility, etc.",
                "RiskSeverity": "Low/Medium/High/Critical",
                "RiskDescription": "Detailed description of the identified risk",
                "Probability": "Likelihood of occurrence (Low/Medium/High)",
                "Impact": "Potential impact on the project (Low/Medium/High)",
                "SecurityImplications": "Any security risks associated",
                "TechnicalMitigation": "Specific technical controls, tools, or implementation details to address the risk",
                "NonTechnicalMitigation": "Process changes, training, policies, and organizational measures to address the risk",
                "ContingencyPlan": "Backup plan in case the risk occurs"
            }}

            IMPORTANT NOTES:
            - Ensure a balanced distribution of risks across all severity levels (Low, Medium, High, Critical).
            - Avoid overestimating severity unless justified by the context.
            - Provide specific, actionable technical and non-technical mitigation strategies.

            Document Section:
            {chunk[:3500]}
            """

            payload = {
                "model": "llama-3.3-70b-versatile",
                "messages": [{"role": "user", "content": prompt}],
                "temperature": 0.4,
                "max_tokens": 1000,
                "response_format": {"type": "json_object"}
            }

            try:
                response = requests.post(
                    "https://api.groq.com/openai/v1/chat/completions",
                    headers=headers,
                    json=payload
                )
                response.raise_for_status()
                content = response.json()["choices"][0]["message"]["content"]
                
                try:
                    json.loads(content)
                    risk_reports.append(content)
                except json.JSONDecodeError as je:
                    logger.error(f"Invalid JSON received for chunk {idx+1}: {str(je)}")
                    risk_reports.append(json.dumps({
                        "RiskID": f"RISK-ERR-{idx+1:03d}",
                        "RiskName": "API Response Parsing Error",
                        "RiskCategory": "Technical",
                        "RiskSeverity": "Low",
                        "RiskDescription": f"The API response for chunk {idx+1} could not be parsed as valid JSON.",
                        "Probability": "Medium",
                        "Impact": "Low",
                        "SecurityImplications": "None",
                        "TechnicalMitigation": "Review the JSON structure and fixing syntax errors in the API integration code",
                        "NonTechnicalMitigation": "Document this parsing issue and establish a review process for analyzing failed responses",
                        "ContingencyPlan": "Contact support if this error persists"
                    }))
                    
                time.sleep(2) 
                break 

            except requests.exceptions.HTTPError as e:
                if response.status_code == 429:  # Too many requests
                    logger.info(f"Rate limited. Retrying in {delay} seconds...")
                    time.sleep(delay)
                    delay *= 2  # Exponential backoff
                    retry_attempts -= 1
                else:
                    logger.error(f"HTTP error: {e}, Status code: {response.status_code}")
                    break  # Stop retrying on non-429 errors
            except Exception as e:
                logger.error(f"Unexpected error: {e}")
                retry_attempts -= 1
                if retry_attempts <= 0:
                    break
    
    return "\n\n".join(risk_reports) if risk_reports else None

def parse_risk_reports(risk_report_text):
    """Parse JSON risk reports and handle formatting issues."""
    risk_items = []
    reports = risk_report_text.split("\n\n")
    
    for idx, report in enumerate(reports):
        try:
            # Try direct JSON parsing first
            try:
                risk_item = json.loads(report)
                risk_items.append(risk_item)
                continue
            except json.JSONDecodeError:
                pass
            
            json_text = None
            
            if "```json" in report:
                json_text = report.split("```json")[1].split("```")[0].strip()
            elif "```" in report:
                json_text = report.split("```")[1].split("```")[0].strip()
            else:
                # Try to find JSON-like content directly
                json_match = re.search(r'(\{.*\})', report, re.DOTALL)
                if json_match:
                    json_text = json_match.group(1).strip()
                else:
                    json_text = report.strip()
            
            if json_text:
                json_text = json_text.replace("'", '"')
                
                json_text = re.sub(r'([{,])\s*(\w+):', r'\1 "\2":', json_text)
                
                json_text = json_text.replace('None', 'null')
                
                risk_item = json.loads(json_text)
                risk_items.append(risk_item)
            else:
                raise ValueError("No JSON content found in report")
                
        except (json.JSONDecodeError, ValueError) as e:
            logger.error(f"Report {idx+1}: Parsing error: {str(e)}")
            risk_items.append({
                "RiskID": f"parsing-error-{idx+1}", 
                "RiskName": "Parsing Error", 
                "RiskDescription": f"Failed to parse JSON: {str(e)}.",
                "RiskSeverity": "Unknown",
                "RiskCategory": "Uncategorized",
                "Probability": "Unknown",
                "Impact": "Unknown",
                "SecurityImplications": "Error parsing this risk report",
                "TechnicalMitigation": "Review the raw text manually",
                "NonTechnicalMitigation": "Establish process for handling parsing errors",
                "ContingencyPlan": "Contact support if multiple parsing errors occur"
            })
        except Exception as e:
            logger.error(f"Report {idx+1}: Processing error: {str(e)}")
            risk_items.append({
                "RiskID": f"processing-error-{idx+1}", 
                "RiskName": "Processing Error", 
                "RiskDescription": f"Error: {str(e)}",
                "RiskSeverity": "Unknown",
                "RiskCategory": "Uncategorized",
                "Probability": "Unknown",
                "Impact": "Unknown",
                "SecurityImplications": "Unknown",
                "TechnicalMitigation": "Review manually",
                "NonTechnicalMitigation": "Document issue",
                "ContingencyPlan": "Retry processing"
            })
            
    return risk_items

# Helper functions
def standardize_severity(severity):
    """Ensure severity is always one of: Critical, High, Medium, Low, Unknown."""
    if not severity:
        return "Unknown"
        
    severity = severity.lower()
    
    if severity in ["red"]:
        return "Critical"
    elif severity in ["orange"]:
        return "High"
    elif severity in ["yellow"]:
        return "Medium"
    elif severity in ["green"]:
        return "Low"
    
    if "critical" in severity:
        return "Critical"
    elif "high" in severity:
        return "High"
    elif "medium" in severity or "med" in severity:
        return "Medium"
    elif "low" in severity:
        return "Low"
    else:
        return "Unknown"

def map_value(value):
    """Map qualitative values to numeric values for RPN calculation."""
    mapping = {"Low": 1, "Medium": 5, "High": 10, "Critical": 15}
    return mapping.get(str(value).strip().title(), 1)  # Default to 1 if unknown

def get_detectability_by_category(category):
    """Determine detectability based on risk category."""
    category = str(category).strip().lower()
    if "security" in category:
        return 9
    elif "technical" in category:
        return 7
    elif "operational" in category:
        return 6
    elif "management" in category:
        return 4
    elif "compliance" in category or "legal" in category:
        return 8
    else:
        return 5 

def map_severity(severity):
    """
    Map qualitative severity to standard FMEA 1-10 scale.
    1 = Minor, 10 = Hazardous without warning
    """
    if not severity or severity == "Unknown":
        return 5  # Default to middle of scale when unknown
        
    severity = str(severity).strip().lower()
    
    # Standard FMEA severity mapping
    if severity in ["critical", "red"]:
        return 9  # Hazardous with warning
    elif severity in ["high", "orange"]:
        return 7  # High severity
    elif severity in ["medium", "med", "yellow"]:
        return 5  # Moderate severity
    elif severity in ["low", "green"]:
        return 3  # Low severity
    else:
        return 5
def map_occurrence(probability):
    """
    Map qualitative probability to standard FMEA 1-10 scale.
    1 = Remote, 10 = Very High
    """
    if not probability or probability == "Unknown":
        return 5  # Default to middle of scale when unknown
        
    probability = str(probability).strip().lower()
    
    # Standard FMEA occurrence mapping
    if probability in ["high", "very high", "certain"]:
        return 8  # High probability
    elif probability in ["medium", "moderate", "likely"]:
        return 5  # Moderate probability  
    elif probability in ["low", "unlikely", "rare"]:
        return 2  # Low probability
    else:
        return 5  # Default to middle of scale

def calculate_detectability(risk):
    """
    Calculate detectability based on risk properties.
    In FMEA, 1 = Almost Certain Detection, 10 = Absolute Uncertainty
    
    This implementation evaluates detectability based on:
    1. Risk category (some risks are inherently harder to detect)
    2. Explicit detection controls mentioned in the risk description
    3. Explicit monitoring systems mentioned in the risk description
    """
    category = str(risk.get("RiskCategory", "")).strip().lower()
    description = str(risk.get("RiskDescription", "")).lower()
    technical_mitigation = str(risk.get("TechnicalMitigation", "")).lower()
    
    # Base detection score by category
    if "security" in category:
        base_score = 7  # Security risks often difficult to detect
    elif "technical" in category:
        base_score = 6  # Technical risks moderately difficult to detect
    elif "operational" in category:
        base_score = 5  # Operational risks moderately detectable
    elif "compliance" in category or "legal" in category:
        base_score = 4  # Compliance risks typically well-defined and detectable
    else:
        base_score = 5  # Default for uncategorized risks
    
    # Adjust for detection controls mentioned
    detection_terms = ["monitor", "alert", "logging", "audit", "detect", "scan", "dashboard", "tracking"]
    control_count = sum(1 for term in detection_terms if term in description or term in technical_mitigation)
    
    # Reduce detection score (easier to detect) if controls are mentioned
    detection_adjustment = min(control_count, 4)  # Cap the adjustment at 4
    
    # Final detection score (clamped between 1-10)
    final_score = max(1, min(10, base_score - detection_adjustment))
    
    return final_score

def extract_current_controls(risk):
    """
    Extract current controls from risk description and mitigations.
    This is a standard FMEA field documenting existing controls.
    """
    description = risk.get("RiskDescription", "")
    technical = risk.get("TechnicalMitigation", "")
    non_technical = risk.get("NonTechnicalMitigation", "")
    
    # Look for mentions of existing controls
    control_indicators = [
        "currently", "existing", "in place", "implemented", 
        "already", "present", "established"
    ]
    
    controls = []
    for text in [description, technical, non_technical]:
        # Extract sentences containing control indicators
        sentences = text.split(". ")
        for sentence in sentences:
            if any(indicator in sentence.lower() for indicator in control_indicators):
                controls.append(sentence.strip())
    
    if not controls:
        return "No existing controls documented."
    
    return " ".join(controls)

def parse_suggested_actions(suggested_fix):
    """
    Parse the AI-generated suggestions into discrete recommended actions.
    """
    if not suggested_fix or suggested_fix == "No immediate action required.":
        return []
        
    # Split by bullet points, numbered lists, or paragraphs
    actions = []
    lines = suggested_fix.split('\n')
    
    current_action = ""
    for line in lines:
        line = line.strip()
        if not line:
            if current_action:
                actions.append(current_action)
                current_action = ""
            continue
            
        # Check if line starts a new action
        if (line.startswith('- ') or line.startswith('• ') or 
            any(line.startswith(f"{i}.") for i in range(1, 11))):
            if current_action:
                actions.append(current_action)
            current_action = line
        else:
            if current_action:
                current_action += " " + line
            else:
                current_action = line
    
    if current_action:
        actions.append(current_action)
        
    # If no structured format found, add the whole text as one action
    if not actions and suggested_fix:
        actions = [suggested_fix]
        
    return actions

def calculate_rpn_and_suggest_fixes(risk_items):
    """
    Calculate RPN for each risk and suggest fixes for high-priority risks.
    Uses standardized FMEA methodology with 1-10 scales.
    """
    high_rpn_threshold = 125  # Standard FMEA often uses 125 as threshold (based on 5×5×5)
    fmea_results = []
    severity_distribution = {"Critical": 0, "High": 0, "Medium": 0, "Low": 0}

    for risk in risk_items:
        try:
            # Map severity to standard 1-10 scale
            severity = map_severity(risk.get("RiskSeverity", "Low"))
            
            occurrence = map_occurrence(risk.get("Probability", "Low"))
            
            detectability = calculate_detectability(risk)

            # Calculate RPN
            rpn = severity * occurrence * detectability
            
            # Add FMEA-specific fields
            risk["FMEA"] = {
                "Severity": severity,
                "Occurrence": occurrence,
                "Detection": detectability,
                "RPN": rpn,
                "CurrentControls": extract_current_controls(risk),
                "RecommendedActions": [],
                "ActionStatus": "Not Started",
                "ResponsiblePerson": "",
                "TargetDate": "",
                "ActionTaken": "",
                "UpdatedRPN": None
            }
            risk["RPN"] = rpn
            # Generate mitigation suggestions for high RPN items
            if rpn >= high_rpn_threshold:
                suggested_actions = generate_ai_suggestions(risk)
                risk["FMEA"]["RecommendedActions"] = parse_suggested_actions(suggested_actions)
                risk["SuggestedFix"] = suggested_actions
            else:
                risk["SuggestedFix"] = "No immediate action required. Monitor as part of regular review."

            # Track severity distribution
            severity_level = risk.get("RiskSeverity", "Unknown")
            if severity_level in severity_distribution:
                severity_distribution[severity_level] += 1

            fmea_results.append(risk)
        except Exception as e:
            logger.error(f"Error calculating RPN for risk: {e}")
            risk["RPN"] = 0
            risk["SuggestedFix"] = f"Error calculating RPN: {e}"
            fmea_results.append(risk)

    logger.info(f"Severity distribution: {severity_distribution}")
    return fmea_results

def generate_ai_suggestions(risk):
    """Generate AI-driven mitigation strategies for high RPN risks."""
    try:
        prompt = f"""
        You are an AI specializing in risk mitigation strategies.
        Given the following risk details, provide specific technical and non-technical mitigation strategies:

        Risk Details:
        - Risk Name: {risk.get('RiskName', 'Unnamed Risk')}
        - Risk Category: {risk.get('RiskCategory', 'Uncategorized')}
        - Risk Severity: {risk.get('RiskSeverity', 'Unknown')}
        - Probability: {risk.get('Probability', 'Unknown')}
        - Impact: {risk.get('Impact', 'Unknown')}
        - Risk Description: {risk.get('RiskDescription', 'No description provided')}

        Provide actionable and detailed mitigation strategies.
        """
        headers = {
            "Authorization": f"Bearer {GROQ_API_KEY}",
            "Content-Type": "application/json"
        }
        payload = {
            "model": "llama-3.3-70b-versatile",
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.4,
            "max_tokens": 500
        }
        response = requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers=headers,
            json=payload
        )
        response.raise_for_status()
        content = response.json()["choices"][0]["message"]["content"]
        return content.strip()
    except Exception as e:
        logger.error(f"Error generating AI suggestions: {e}")
        return "Error generating AI suggestions. Please review the risk manually."

def calculate_overall_risk(risk_items):
    """Calculate the overall risk level and generate a summary."""
    risk_levels = {"Critical": 4, "High": 3, "Medium": 2, "Low": 1}
    max_risk_level = 0
    summary = []

    for risk in risk_items:
        severity = risk.get("RiskSeverity", "Low")
        max_risk_level = max(max_risk_level, risk_levels.get(severity, 1))
        summary.append(f"{risk.get('RiskName', 'Unnamed Risk')}: {severity}")

    overall_level = next((key for key, value in risk_levels.items() if value == max_risk_level), "Low")
    return overall_level, summary

@app.route('/api/upload', methods=['POST'])
def upload_file():
    """API endpoint to handle file uploads."""
    if 'file' not in request.files:
        return jsonify({"error": "No file part in the request"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No file selected"}), 400
    
    if not file.filename.lower().endswith(('.pdf', '.docx', '.txt', '.ppt', '.pptx')):
        return jsonify({"error": "Unsupported file format. Please upload a PDF, DOCX, TXT, or PPT/PPTX file."}), 400
    
    # Save file
    filename = secure_filename(file.filename)
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(file_path)
    
    file_extension = filename.split('.')[-1].lower()
    extracted_text = None
    
    if file_extension == "pdf":
        extracted_text = extract_text_from_pdf(file_path)
    elif file_extension == "docx":
        extracted_text = extract_text_from_docx(file_path)
    elif file_extension == "txt":
        extracted_text = extract_text_from_txt(file_path)
    elif file_extension in ["ppt", "pptx"]:
        extracted_text = extract_text_from_pptx(file_path)
    
    os.remove(file_path)
    
    if not extracted_text:
        return jsonify({"error": "Failed to extract text from the document"}), 500
    
    user_id = request.headers.get('User-ID')
    if not user_id:
        logger.error("User ID is missing in the request headers.")
        logger.info(f"Request headers: {request.headers}")
        return jsonify({"error": "User ID is required to associate the upload with a user."}), 400

    logger.info(f"Received User ID: {user_id}")

    risk_report = analyze_risks_with_groq(extracted_text)
    if not risk_report:
        return jsonify({"error": "Failed to generate risk assessment report"}), 500
    
    risk_items = parse_risk_reports(risk_report)
    
    risk_items = calculate_rpn_and_suggest_fixes(risk_items)
    
    for item in risk_items:
        if "RiskSeverity" in item:
            item["RiskSeverity"] = standardize_severity(item["RiskSeverity"])
    
    overall_level, summary = calculate_overall_risk(risk_items)

    history_entry = {
        "user_id": user_id,
        "file_name": filename,
        "description": f"Uploaded {filename} for risk analysis.",
        "upload_date": datetime.datetime.now(),
        "risk_summary": {
            "level": overall_level,
            "summary": ", ".join(summary),
            "details": risk_items
        }
    }
    history_collection.insert_one(history_entry)

    return jsonify({"success": True, "risk_items": risk_items})

@app.route('/api/history', methods=['GET'])
def get_user_history():
    user_id = request.headers.get('User-ID') 
    if not user_id:
        return jsonify({"error": "User ID is required"}), 400

    history = history_collection.find({"user_id": user_id})
    history_data = [
        {
            "file_name": entry.get("file_name", ""),
            "description": entry.get("description", ""),
            "upload_date": entry.get("upload_date", ""),
            "risk_summary": entry.get("risk_summary", {})
        }
        for entry in history
    ]
    return jsonify({"success": True, "history": history_data})

@app.route('/api/history', methods=['DELETE'])
def delete_history_item():
    user_id = request.headers.get('User-ID')
    file_name = request.json.get('file_name')
    if not user_id or not file_name:
        return jsonify({"error": "User ID and file name are required"}), 400

    result = history_collection.delete_one({"user_id": user_id, "file_name": file_name})
    if result.deleted_count == 1:
        return jsonify({"success": True})
    else:
        return jsonify({"success": False, "error": "Document not found"}), 404

@app.route('/api/health', methods=['GET'])
def health_check():
    """API endpoint for health check."""
    return jsonify({"status": "alive"})

if __name__ == "__main__":
    app.run(debug=True, port=5001)